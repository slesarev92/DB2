"""Integration tests для PPTX экспорта (5.2).

Покрывает:
- Service-level: `generate_project_pptx → bytes` + парсится через
  python-pptx, 13 слайдов, ключевые заголовки на месте, content fields
  попадают в слайды
- Endpoint-level: GET /api/projects/{id}/export/pptx → 200 + правильный
  MIME + filename, 404, 401
"""
from __future__ import annotations

from io import BytesIO

import pytest
from httpx import AsyncClient
from pptx import Presentation
from sqlalchemy.ext.asyncio import AsyncSession

from app.export.excel_exporter import ProjectNotFoundForExport
from app.export.ppt_exporter import generate_project_pptx
from tests.api.test_calculation import _seed_minimal_project


PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)


def _collect_slide_text(prs: Presentation) -> list[str]:
    """Собирает весь текст со всех слайдов — включая ячейки таблиц.

    Таблицы в python-pptx это GraphicFrame shapes с `has_table`.
    Без обхода table.cells значительная часть контента теряется
    (KPI, PnL, validation, function_readiness, roadmap, approvers).
    """
    out: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        if r.text:
                            out.append(r.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            for r in p.runs:
                                if r.text:
                                    out.append(r.text)
    return out


def _get_slide_title(slide) -> str:
    """Первый непустой текст слайда — обычно это title."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip():
                        return r.text.strip()
    return ""


# ============================================================
# Service-level tests
# ============================================================


class TestGenerateProjectPptx:
    async def test_returns_bytes(self, db_session: AsyncSession):
        project_id, _, _, _ = await _seed_minimal_project(db_session)
        result = await generate_project_pptx(db_session, project_id)
        assert isinstance(result, bytes)
        assert len(result) > 0
        # ZIP signature (PPTX — это ZIP)
        assert result[:2] == b"PK"

    async def test_has_13_slides(self, db_session: AsyncSession):
        project_id, _, _, _ = await _seed_minimal_project(db_session)
        result = await generate_project_pptx(db_session, project_id)

        prs = Presentation(BytesIO(result))
        assert len(prs.slides) == 13

    async def test_slide_titles(self, db_session: AsyncSession):
        """Каждый слайд имеет ожидаемый заголовок."""
        project_id, _, _, _ = await _seed_minimal_project(db_session)
        result = await generate_project_pptx(db_session, project_id)

        prs = Presentation(BytesIO(result))
        slides = list(prs.slides)

        # Слайд 1 — название проекта (из _seed_minimal_project)
        assert "Calc test project" in _get_slide_title(slides[0])

        expected_titles = [
            "1. Общая информация",
            "2. Концепция продукта",
            "3. Технология и обоснование",
            "4. Результаты валидации",
            "5. Продуктовый микс",
            "6. Финансовая модель",
            "7. Ключевые KPI",
            "8. PnL по годам",
            "9. Стакан себестоимости",
            "10. Риски и готовность функций",
            "11. Дорожная карта и согласующие",
            "12. Executive Summary",
        ]
        for i, expected in enumerate(expected_titles, start=1):
            title = _get_slide_title(slides[i])
            assert expected in title, (
                f"Слайд {i + 1} ожидался «{expected}», получили «{title}»"
            )

    async def test_content_fields_appear_in_slides(
        self, db_session: AsyncSession
    ):
        """Если content поля заполнены — они попадают в PPTX."""
        from datetime import date as _date

        from app.models import Project

        project_id, _, _, _ = await _seed_minimal_project(db_session)
        project = await db_session.get(Project, project_id)
        assert project is not None
        # Заполняем несколько content полей (Фаза 4.5)
        project.description = "Уникальная спортивная вода с электролитами"
        project.gate_stage = "G3"
        project.passport_date = _date(2025, 9, 1)
        project.project_owner = "Иван Иванов"
        project.growth_opportunity = "Сегмент спорт-воды растёт 15% в год"
        project.risks = ["Конкуренты", "Регулятор"]
        project.function_readiness = {
            "R&D": {"status": "green", "notes": "рецепт финализирован"},
            "Marketing": {"status": "yellow", "notes": "комстрат в работе"},
        }
        project.roadmap_tasks = [
            {
                "name": "Первая партия",
                "start_date": "2025-04-01",
                "status": "in_progress",
                "owner": "Production",
            }
        ]
        project.approvers = [
            {"metric": "NPV", "name": "CFO", "source": "модель"}
        ]
        project.executive_summary = "Проект окупится за 4 года при Base сценарии."
        await db_session.flush()

        result = await generate_project_pptx(db_session, project_id)
        prs = Presentation(BytesIO(result))
        all_text = " ".join(_collect_slide_text(prs))

        # Scalar поля
        assert "спортивная вода с электролитами" in all_text
        assert "Иван Иванов" in all_text
        assert "G3" in all_text
        assert "Сегмент спорт-воды растёт" in all_text
        assert "Проект окупится за 4 года" in all_text

        # JSONB — риски
        assert "Конкуренты" in all_text
        # function_readiness
        assert "R&D" in all_text
        assert "рецепт финализирован" in all_text
        # roadmap
        assert "Первая партия" in all_text
        # approvers
        assert "CFO" in all_text

    async def test_handles_empty_content_fields(
        self, db_session: AsyncSession
    ):
        """Проект без content полей → PPTX генерируется, секции с «—»."""
        project_id, _, _, _ = await _seed_minimal_project(db_session)
        # Не трогаем content — все null после _seed_minimal_project

        result = await generate_project_pptx(db_session, project_id)
        prs = Presentation(BytesIO(result))
        assert len(prs.slides) == 13

        # Placeholder'ы должны быть хотя бы в одной ячейке каждой секции
        all_text = " ".join(_collect_slide_text(prs))
        assert "—" in all_text

    async def test_404_for_unknown_project(self, db_session: AsyncSession):
        with pytest.raises(ProjectNotFoundForExport):
            await generate_project_pptx(db_session, 999999)


# ============================================================
# Endpoint tests
# ============================================================


class TestExportPptxEndpoint:
    async def test_returns_pptx_with_correct_mime(
        self,
        auth_client: AsyncClient,
        db_session: AsyncSession,
    ):
        project_id, _, _, _ = await _seed_minimal_project(db_session)
        await db_session.commit()

        resp = await auth_client.get(
            f"/api/projects/{project_id}/export/pptx"
        )
        assert resp.status_code == 200, resp.text
        assert resp.headers["content-type"].startswith(PPTX_MIME)

    async def test_response_is_valid_pptx(
        self,
        auth_client: AsyncClient,
        db_session: AsyncSession,
    ):
        project_id, _, _, _ = await _seed_minimal_project(db_session)
        await db_session.commit()

        resp = await auth_client.get(
            f"/api/projects/{project_id}/export/pptx"
        )
        assert resp.status_code == 200

        prs = Presentation(BytesIO(resp.content))
        assert len(prs.slides) == 13

    async def test_filename_in_content_disposition(
        self,
        auth_client: AsyncClient,
        db_session: AsyncSession,
    ):
        project_id, _, _, _ = await _seed_minimal_project(db_session)
        await db_session.commit()

        resp = await auth_client.get(
            f"/api/projects/{project_id}/export/pptx"
        )
        cd = resp.headers.get("content-disposition", "")
        assert "attachment" in cd
        assert f"project_{project_id}" in cd
        assert ".pptx" in cd

    async def test_404_for_unknown_project(
        self, auth_client: AsyncClient
    ):
        resp = await auth_client.get("/api/projects/999999/export/pptx")
        assert resp.status_code == 404

    async def test_unauthorized(self, client: AsyncClient):
        resp = await client.get("/api/projects/1/export/pptx")
        assert resp.status_code == 401

    async def test_cyrillic_project_name_does_not_break_header(
        self,
        auth_client: AsyncClient,
        db_session: AsyncSession,
    ):
        """Regression: проект с кириллическим названием должен экспортиться.

        HTTP headers — latin-1 only (RFC 7230). Прямое filename="{name}"
        с русским именем даёт UnicodeEncodeError. Исправлено через RFC
        5987 filename*=UTF-8''{percent} + ASCII fallback в filename=.
        """
        from app.models import Project

        project_id, _, _, _ = await _seed_minimal_project(db_session)
        project = await db_session.get(Project, project_id)
        assert project is not None
        project.name = "Тестовый проект с кириллицей"
        await db_session.flush()
        await db_session.commit()

        resp = await auth_client.get(
            f"/api/projects/{project_id}/export/pptx"
        )
        assert resp.status_code == 200, resp.text
        cd = resp.headers.get("content-disposition", "")
        # ASCII fallback в filename=
        assert "filename=" in cd
        # RFC 5987 UTF-8 версия — кириллица percent-encoded
        assert "filename*=UTF-8''" in cd
        # %D0 — первый байт кириллицы в UTF-8
        assert "%D0" in cd
