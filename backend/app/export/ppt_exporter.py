"""Генератор PPTX-экспорта паспорта проекта (задача 5.2, F-09).

Структура презентации — 13 слайдов, покрывающих 4 блока:

Content (Фаза 4.5):
  1.  Титульный: название проекта, brand, gate_stage, passport_date, owner
  2.  Общая информация: description, project_goal, geography, innovation_type
  3.  Концепция продукта: growth_opportunity, concept_text, idea_short,
      target_audience, replacement_target
  4.  Технология + R&D: technology, rnd_progress, rationale
  5.  Результаты валидации: 5 подтестов concept_test/naming/design/product/price
  6.  Продуктовый микс: SKU таблица + package images

Financials:
  7.  Финансовая модель: WACC / Tax / WC / VAT / горизонт / инфляция
  8.  KPI: NPV/IRR/ROI/Payback × 3 сценария × Y1-Y10
  9.  PnL по годам: Y1..Y10 агрегаты (из base pipeline)
  10. Стакан себестоимости: BOM top items per SKU + финансовый план CAPEX/OPEX

Roadmap & governance:
  11. Риски + Готовность функций (2-колонный слайд)
  12. Дорожная карта + Согласующие
  13. Executive summary

Реализация через `python-pptx` (MIT). Только стандартные layouts —
никакого corporate template для MVP. Размер слайда 16:9 (default).

Контракт как у excel_exporter:
- `generate_project_pptx(session, project_id) → bytes`
- Загрузка данных переиспользует helpers из `excel_exporter` (не DRY
  violation — оба модуля служат одному use case «экспорт проекта»).
- Если расчёт не выполнен — KPI/PnL слайды остаются с пометкой
  «Расчёт не выполнен».
- Если content поля null — секции остаются с «—».
"""
from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.util import Emu, Inches, Pt
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.core.config import settings
from app.engine.pipeline import run_project_pipeline
from app.export.excel_exporter import (
    ProjectNotFoundForExport,
    _load_project_full,
    _load_psk_channels,
    _load_scenario_results,
    _load_skus_with_bom,
)
from app.models import (
    BOMItem,
    MediaAsset,
    Project,
    ProjectFinancialPlan,
    ProjectSKU,
    ProjectSKUChannel,
    RefInflation,
    Scenario,
    ScenarioResult,
    ScenarioType,
)
from app.models.base import PeriodScope
from app.services.calculation_service import (
    _load_period_catalog,
    _load_project_financial_plan,
    build_line_inputs,
)


# ============================================================
# Layout constants (16:9 default: 13.333" × 7.5")
# ============================================================

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Стандартные layouts python-pptx:
#   0 Title Slide
#   1 Title and Content
#   5 Title Only
#   6 Blank
LAYOUT_TITLE = 0
LAYOUT_TITLE_AND_CONTENT = 1
LAYOUT_TITLE_ONLY = 5
LAYOUT_BLANK = 6

TITLE_FONT_SIZE = Pt(28)
SECTION_FONT_SIZE = Pt(18)
BODY_FONT_SIZE = Pt(12)
SMALL_FONT_SIZE = Pt(10)

# Коды стадий → русская подпись (для слайда 1)
GATE_LABELS = {
    "G0": "G0 — Идея",
    "G1": "G1 — Концепция",
    "G2": "G2 — Design",
    "G3": "G3 — Development",
    "G4": "G4 — Launch Ready",
    "G5": "G5 — In Market",
}

SCENARIO_LABELS = {
    ScenarioType.BASE: "Base",
    ScenarioType.CONSERVATIVE: "Conservative",
    ScenarioType.AGGRESSIVE: "Aggressive",
}

FUNCTION_STATUS_LABELS = {
    "green": "🟢 Готово",
    "yellow": "🟡 В работе",
    "red": "🔴 Риск",
}


# ============================================================
# Helpers
# ============================================================


def _fmt_money(value: float | int | None, decimals: int = 0) -> str:
    """Форматирует денежное значение с разделителями. None → «—»."""
    if value is None:
        return "—"
    try:
        f = float(value)
    except (TypeError, ValueError):
        return "—"
    if decimals == 0:
        return f"{f:,.0f}".replace(",", " ")
    return f"{f:,.{decimals}f}".replace(",", " ")


def _fmt_pct(value: float | int | None, decimals: int = 1) -> str:
    if value is None:
        return "—"
    try:
        return f"{float(value) * 100:.{decimals}f}%"
    except (TypeError, ValueError):
        return "—"


def _fmt_text(value: str | None) -> str:
    """None / empty → «—»."""
    return value.strip() if value and value.strip() else "—"


def _add_text_box(
    slide: Slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    text: str,
    *,
    bold: bool = False,
    size: Pt | None = None,
    align: int | None = None,
) -> None:
    """Добавляет прямоугольник с текстом в слайд."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    if bold:
        run.font.bold = True
    if size is not None:
        run.font.size = size
    if align is not None:
        p.alignment = align


def _add_title(slide: Slide, title_text: str) -> None:
    """Добавляет стандартный заголовок слайда (используется с layouts,
    у которых есть title placeholder). Если placeholder'а нет, добавляет
    textbox сверху.
    """
    if slide.shapes.title is not None:
        slide.shapes.title.text = title_text
        for p in slide.shapes.title.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(24)
                r.font.bold = True
    else:
        _add_text_box(
            slide,
            Inches(0.5),
            Inches(0.3),
            Inches(12.3),
            Inches(0.6),
            title_text,
            bold=True,
            size=Pt(24),
        )


def _add_field_block(
    slide: Slide,
    left: Emu,
    top: Emu,
    width: Emu,
    label: str,
    value: str,
    *,
    value_height: Emu = Inches(0.4),
) -> Emu:
    """Рисует блок «Label: value» (label bold sm + value reg md).

    Возвращает новую вертикальную позицию (top + занятая высота).
    """
    # Label
    _add_text_box(
        slide,
        left,
        top,
        width,
        Inches(0.25),
        label,
        bold=True,
        size=SMALL_FONT_SIZE,
    )
    # Value
    _add_text_box(
        slide,
        left,
        top + Inches(0.25),
        width,
        value_height,
        value,
        size=BODY_FONT_SIZE,
    )
    return top + Inches(0.25) + value_height + Inches(0.15)


def _add_simple_table(
    slide: Slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    headers: list[str],
    rows: list[list[str]],
) -> None:
    """Вставляет простую таблицу без cell merging.

    Если rows пустой — таблица всё равно вставится (только заголовок).
    """
    rows_count = max(len(rows), 1) + 1  # +1 для header
    cols_count = len(headers)
    table_shape = slide.shapes.add_table(
        rows_count, cols_count, left, top, width, height
    )
    table = table_shape.table

    # Заголовки
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = SMALL_FONT_SIZE

    # Данные
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            if c_idx >= cols_count:
                break
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = SMALL_FONT_SIZE

    # Если rows пустой — пишем "—" в placeholder-строке
    if not rows:
        for c in range(cols_count):
            cell = table.cell(1, c)
            cell.text = "—"


def _add_bullets(
    slide: Slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    items: list[str],
    *,
    placeholder: str = "Нет данных",
) -> None:
    """Bullet-список. Если items пустой — пишет placeholder."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if not items:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = placeholder
        run.font.size = BODY_FONT_SIZE
        run.font.italic = True
        return
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = BODY_FONT_SIZE


# ============================================================
# Slide builders
# ============================================================


def _build_slide_title(prs: PresentationType, project: Project) -> None:
    """Слайд 1: титульный."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])

    # Большой заголовок с именем проекта
    _add_text_box(
        slide,
        Inches(0.5),
        Inches(2.5),
        Inches(12.3),
        Inches(1.5),
        project.name,
        bold=True,
        size=Pt(44),
        align=PP_ALIGN.CENTER,
    )

    # Подзаголовок
    subtitle_parts: list[str] = ["Цифровой паспорт проекта"]
    if project.gate_stage:
        subtitle_parts.append(GATE_LABELS.get(project.gate_stage, project.gate_stage))
    _add_text_box(
        slide,
        Inches(0.5),
        Inches(4.2),
        Inches(12.3),
        Inches(0.6),
        " · ".join(subtitle_parts),
        size=Pt(20),
        align=PP_ALIGN.CENTER,
    )

    # Phase 8.7: Gate timeline (текстовая визуализация G0 → G5)
    gate_str = _build_gate_timeline_text(project.gate_stage)
    _add_text_box(
        slide,
        Inches(0.5),
        Inches(5.0),
        Inches(12.3),
        Inches(0.6),
        gate_str,
        size=Pt(16),
        align=PP_ALIGN.CENTER,
    )

    # Footer: owner + passport_date
    footer_parts: list[str] = []
    if project.project_owner:
        footer_parts.append(f"Владелец: {project.project_owner}")
    if project.passport_date:
        footer_parts.append(f"Дата паспорта: {project.passport_date.isoformat()}")
    footer_parts.append(f"Старт проекта: {project.start_date.isoformat()}")
    _add_text_box(
        slide,
        Inches(0.5),
        Inches(6.5),
        Inches(12.3),
        Inches(0.5),
        " · ".join(footer_parts),
        size=Pt(14),
        align=PP_ALIGN.CENTER,
    )


def _build_gate_timeline_text(current_gate: str | None) -> str:
    """Phase 8.7: Текстовая визуализация G0 → G5 с маркером текущего гейта."""
    gates = ["G0", "G1", "G2", "G3", "G4", "G5"]
    parts: list[str] = []
    for g in gates:
        if g == current_gate:
            parts.append(f"●{g}")  # текущий
        else:
            parts.append(g)
    return "  →  ".join(parts)


def _build_slide_general_info(prs: PresentationType, project: Project) -> None:
    """Слайд 2: общая информация."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _add_title(slide, "1. Общая информация")

    # Левая колонка: 4 поля
    left_x = Inches(0.5)
    right_x = Inches(6.9)
    col_width = Inches(6.0)
    y = Inches(1.2)
    y = _add_field_block(
        slide, left_x, y, col_width,
        "Описание проекта", _fmt_text(project.description),
        value_height=Inches(1.0),
    )
    y = _add_field_block(
        slide, left_x, y, col_width,
        "Цель проекта", _fmt_text(project.project_goal),
        value_height=Inches(1.0),
    )
    y = _add_field_block(
        slide, left_x, y, col_width,
        "География", _fmt_text(project.geography),
    )
    y = _add_field_block(
        slide, left_x, y, col_width,
        "Владелец проекта", _fmt_text(project.project_owner),
    )

    # Правая колонка: 4 поля
    y = Inches(1.2)
    y = _add_field_block(
        slide, right_x, y, col_width,
        "Gate-стадия",
        GATE_LABELS.get(project.gate_stage, "—") if project.gate_stage else "—",
    )
    y = _add_field_block(
        slide, right_x, y, col_width,
        "Дата паспорта",
        project.passport_date.isoformat() if project.passport_date else "—",
    )
    y = _add_field_block(
        slide, right_x, y, col_width,
        "Тип инновации", _fmt_text(project.innovation_type),
    )
    y = _add_field_block(
        slide, right_x, y, col_width,
        "Тип производства", _fmt_text(project.production_type),
    )


def _build_slide_concept(prs: PresentationType, project: Project) -> None:
    """Слайд 3: концепция продукта."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _add_title(slide, "2. Концепция продукта")

    y = Inches(1.2)
    left = Inches(0.5)
    width = Inches(12.3)

    y = _add_field_block(
        slide, left, y, width,
        "Growth opportunity", _fmt_text(project.growth_opportunity),
        value_height=Inches(0.6),
    )
    y = _add_field_block(
        slide, left, y, width,
        "Концепция", _fmt_text(project.concept_text),
        value_height=Inches(1.0),
    )
    y = _add_field_block(
        slide, left, y, width,
        "Короткая идея", _fmt_text(project.idea_short),
    )
    y = _add_field_block(
        slide, left, y, width,
        "Целевая аудитория", _fmt_text(project.target_audience),
        value_height=Inches(0.6),
    )
    y = _add_field_block(
        slide, left, y, width,
        "Кого замещаем", _fmt_text(project.replacement_target),
    )


def _build_slide_technology(prs: PresentationType, project: Project) -> None:
    """Слайд 4: технология, R&D, обоснование."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _add_title(slide, "3. Технология и обоснование")

    y = Inches(1.2)
    left = Inches(0.5)
    width = Inches(12.3)

    y = _add_field_block(
        slide, left, y, width,
        "Технология", _fmt_text(project.technology),
        value_height=Inches(1.0),
    )
    y = _add_field_block(
        slide, left, y, width,
        "Прогресс R&D", _fmt_text(project.rnd_progress),
        value_height=Inches(1.0),
    )
    y = _add_field_block(
        slide, left, y, width,
        "Обоснование (rationale)", _fmt_text(project.rationale),
        value_height=Inches(1.5),
    )


def _build_slide_validation(prs: PresentationType, project: Project) -> None:
    """Слайд 5: результаты валидации (5 подтестов)."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _add_title(slide, "4. Результаты валидации")

    subtests = [
        ("concept_test", "Concept test"),
        ("naming", "Naming"),
        ("design", "Design"),
        ("product", "Product"),
        ("price", "Price"),
    ]
    rows: list[list[str]] = []
    validation = project.validation_tests or {}
    for key, label in subtests:
        entry = validation.get(key) if isinstance(validation, dict) else None
        if not isinstance(entry, dict):
            entry = {}
        score = entry.get("score")
        notes = entry.get("notes", "") or ""
        score_str = "—" if score is None else str(score)
        rows.append([label, score_str, notes])

    _add_simple_table(
        slide,
        Inches(0.5),
        Inches(1.2),
        Inches(12.3),
        Inches(5.5),
        ["Подтест", "Score", "Notes"],
        rows,
    )


def _build_slide_sku_mix(
    prs: PresentationType,
    skus_with_bom: list[tuple[ProjectSKU, list[BOMItem]]],
    package_images: dict[int, Path],
) -> None:
    """Слайд 6: продуктовый микс (таблица SKU + package images).

    package_images: dict {media_id: local Path} — файлы на диске для
    embedding'а в слайд через slide.shapes.add_picture.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _add_title(slide, "5. Продуктовый микс")

    # Таблица SKU сверху
    rows: list[list[str]] = []
    for psk, bom_items in skus_with_bom:
        bom_cost = sum(
            float(b.quantity_per_unit) * float(b.price_per_unit) * (1 + float(b.loss_pct))
            for b in bom_items
        )
        rows.append([
            f"{psk.sku.brand} — {psk.sku.name}",
            psk.sku.format or "—",
            f"{float(psk.sku.volume_l):.2f} л" if psk.sku.volume_l else "—",
            _fmt_pct(float(psk.production_cost_rate)),
            _fmt_money(bom_cost, decimals=2),
        ])

    _add_simple_table(
        slide,
        Inches(0.5),
        Inches(1.2),
        Inches(12.3),
        Inches(3.0),
        ["SKU", "Формат", "Объём", "Prod cost rate", "BOM ₽/ед"],
        rows,
    )

    # Package images ниже (до 4 штук в ряд)
    img_top = Inches(4.6)
    img_size = Inches(1.8)
    img_gap = Inches(0.3)
    img_idx = 0
    max_images = 6
    for psk, _bom in skus_with_bom:
        if psk.package_image_id is None:
            continue
        img_path = package_images.get(psk.package_image_id)
        if img_path is None or not img_path.is_file():
            continue
        if img_idx >= max_images:
            break
        x = Inches(0.5) + img_idx * (img_size + img_gap)
        try:
            slide.shapes.add_picture(
                str(img_path), x, img_top, width=img_size, height=img_size
            )
        except Exception:  # noqa: BLE001
            # Битый PNG или неподдерживаемый формат — пропускаем,
            # презентация всё равно сгенерируется.
            continue
        # Подпись под картинкой
        _add_text_box(
            slide,
            x,
            img_top + img_size + Inches(0.05),
            img_size,
            Inches(0.3),
            f"{psk.sku.brand} — {psk.sku.name}",
            size=SMALL_FONT_SIZE,
            align=PP_ALIGN.CENTER,
        )
        img_idx += 1

    if img_idx == 0:
        _add_text_box(
            slide,
            Inches(0.5),
            img_top,
            Inches(12.3),
            Inches(0.4),
            "Изображения упаковки не загружены",
            size=SMALL_FONT_SIZE,
        )


def _build_slide_financial_model(
    prs: PresentationType,
    project: Project,
    inflation_profile: RefInflation | None,
) -> None:
    """Слайд 7: макро-факторы финансовой модели."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _add_title(slide, "6. Финансовая модель — макро-факторы")

    params = [
        ["WACC", _fmt_pct(float(project.wacc), decimals=2)],
        ["Налог на прибыль", _fmt_pct(float(project.tax_rate), decimals=2)],
        ["Working Capital rate", _fmt_pct(float(project.wc_rate), decimals=2)],
        ["VAT rate", _fmt_pct(float(project.vat_rate), decimals=2)],
        ["Валюта", project.currency],
        ["Горизонт", f"{project.horizon_years} лет"],
        [
            "Профиль инфляции",
            inflation_profile.profile_name if inflation_profile else "—",
        ],
    ]
    _add_simple_table(
        slide,
        Inches(0.5),
        Inches(1.2),
        Inches(12.3),
        Inches(4.0),
        ["Параметр", "Значение"],
        params,
    )


def _build_slide_kpi(
    prs: PresentationType,
    scenarios: list[Scenario],
    results_by_scenario: dict[int, list[ScenarioResult]],
) -> None:
    """Слайд 8: сводная KPI таблица 3 сценария × Y1-Y10."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _add_title(slide, "7. Ключевые KPI — сравнение сценариев")

    scenario_order = {
        ScenarioType.BASE: 0,
        ScenarioType.CONSERVATIVE: 1,
        ScenarioType.AGGRESSIVE: 2,
    }
    sorted_scenarios = sorted(scenarios, key=lambda s: scenario_order.get(s.type, 99))

    rows: list[list[str]] = []
    for sc in sorted_scenarios:
        results = results_by_scenario.get(sc.id, [])
        y1y10 = next(
            (r for r in results if r.period_scope == PeriodScope.Y1Y10), None
        )
        if y1y10 is None:
            rows.append([
                SCENARIO_LABELS.get(sc.type, sc.type.value),
                "—", "—", "—", "—", "—",
            ])
        else:
            rows.append([
                SCENARIO_LABELS.get(sc.type, sc.type.value),
                _fmt_money(float(y1y10.npv) if y1y10.npv is not None else None),
                _fmt_pct(float(y1y10.irr) if y1y10.irr is not None else None),
                _fmt_pct(float(y1y10.roi) if y1y10.roi is not None else None),
                (
                    f"{float(y1y10.payback_simple):.1f}"
                    if y1y10.payback_simple is not None
                    else "—"
                ),
                "✓" if y1y10.go_no_go else ("✗" if y1y10.go_no_go is False else "—"),
            ])

    _add_simple_table(
        slide,
        Inches(0.5),
        Inches(1.4),
        Inches(12.3),
        Inches(2.0),
        ["Сценарий", "NPV ₽", "IRR", "ROI", "Payback, лет", "Go/No-Go"],
        rows,
    )

    # Phase 8.3: Per-unit метрики Base сценария (если посчитан)
    base_scenario = next(
        (s for s in sorted_scenarios if s.type == ScenarioType.BASE), None
    )
    if base_scenario is not None:
        base_results = results_by_scenario.get(base_scenario.id, [])
        scope_order = [PeriodScope.Y1Y3, PeriodScope.Y1Y5, PeriodScope.Y1Y10]
        scope_labels = ["Y1-Y3", "Y1-Y5", "Y1-Y10"]
        per_unit_rows: list[list[str]] = []
        for metric_name, attr in [
            ("Выручка / шт, ₽", "nr_per_unit"),
            ("GP / шт, ₽", "gp_per_unit"),
            ("CM / шт, ₽", "cm_per_unit"),
            ("EBITDA / шт, ₽", "ebitda_per_unit"),
        ]:
            row = [metric_name]
            for scope in scope_order:
                r = next((x for x in base_results if x.period_scope == scope), None)
                val = getattr(r, attr, None) if r else None
                row.append(
                    _fmt_money(float(val), 2) if val is not None else "—"
                )
            per_unit_rows.append(row)

        _add_text_box(
            slide,
            Inches(0.5),
            Inches(3.6),
            Inches(12.3),
            Inches(0.4),
            "Per-unit экономика (Base сценарий, scope-средняя):",
            bold=True,
            size=SMALL_FONT_SIZE,
        )
        _add_simple_table(
            slide,
            Inches(0.5),
            Inches(4.0),
            Inches(12.3),
            Inches(2.5),
            ["Метрика"] + scope_labels,
            per_unit_rows,
        )


def _build_slide_pnl(
    prs: PresentationType,
    base_aggregate: Any | None,
) -> None:
    """Слайд 9: PnL по годам Y1..Y10 (из base pipeline agregate)."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _add_title(slide, "8. PnL по годам — Base сценарий")

    if base_aggregate is None or not base_aggregate.annual_free_cash_flow:
        _add_text_box(
            slide,
            Inches(0.5),
            Inches(3.0),
            Inches(12.3),
            Inches(0.5),
            "Расчёт не выполнен. Запустите пересчёт проекта перед экспортом.",
            size=BODY_FONT_SIZE,
            align=PP_ALIGN.CENTER,
        )
        return

    years = list(range(1, len(base_aggregate.annual_free_cash_flow) + 1))
    headers = ["Метрика"] + [f"Y{y}" for y in years]

    metric_rows: list[list[str]] = []
    annual_metrics = [
        ("Net Revenue", base_aggregate.annual_net_revenue),
        ("Contribution", base_aggregate.annual_contribution),
        ("FCF", base_aggregate.annual_free_cash_flow),
        ("DCF", base_aggregate.annual_discounted_cash_flow),
        ("Cumulative FCF", base_aggregate.cumulative_fcf),
    ]
    for label, values in annual_metrics:
        row = [label]
        for v in values:
            row.append(_fmt_money(v))
        metric_rows.append(row)

    _add_simple_table(
        slide,
        Inches(0.3),
        Inches(1.2),
        Inches(12.7),
        Inches(5.5),
        headers,
        metric_rows,
    )


def _build_slide_cogs_and_fin_plan(
    prs: PresentationType,
    skus_with_bom: list[tuple[ProjectSKU, list[BOMItem]]],
    financial_plan: list[ProjectFinancialPlan],
    period_by_id: dict[int, Any],
    opex_by_category: dict[str, float] | None = None,
) -> None:
    """Слайд 10: стакан себестоимости + financial plan (2 колонки)."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _add_title(slide, "9. Стакан себестоимости и финансовый план")

    # Левая колонка: BOM топ
    _add_text_box(
        slide,
        Inches(0.5),
        Inches(1.2),
        Inches(6.0),
        Inches(0.4),
        "Стакан BOM (топ ингредиентов)",
        bold=True,
        size=Pt(14),
    )

    bom_rows: list[list[str]] = []
    # Агрегируем топ-10 по всем SKU
    ingredient_costs: dict[str, float] = {}
    for _psk, bom_items in skus_with_bom:
        for b in bom_items:
            cost = (
                float(b.quantity_per_unit)
                * float(b.price_per_unit)
                * (1 + float(b.loss_pct))
            )
            ingredient_costs[b.ingredient_name] = (
                ingredient_costs.get(b.ingredient_name, 0.0) + cost
            )
    top_ingredients = sorted(
        ingredient_costs.items(), key=lambda x: x[1], reverse=True
    )[:10]
    for name, cost in top_ingredients:
        bom_rows.append([name, _fmt_money(cost, decimals=2)])

    _add_simple_table(
        slide,
        Inches(0.5),
        Inches(1.7),
        Inches(6.0),
        Inches(5.0),
        ["Ингредиент", "Стоимость ₽/ед"],
        bom_rows,
    )

    # Правая колонка: financial plan
    _add_text_box(
        slide,
        Inches(6.8),
        Inches(1.2),
        Inches(6.0),
        Inches(0.4),
        "Финансовый план (CAPEX / OPEX по годам)",
        bold=True,
        size=Pt(14),
    )
    # Агрегируем CAPEX/OPEX по годам через model_year периода.
    # Plan хранится per period_id, не per year — один период обычно
    # покрывает один model_year, но суммируем на случай нескольких.
    fp_by_year: dict[int, tuple[float, float]] = {}
    for fp in financial_plan:
        period = period_by_id.get(fp.period_id)
        if period is None:
            continue
        year = period.model_year
        prev_capex, prev_opex = fp_by_year.get(year, (0.0, 0.0))
        fp_by_year[year] = (
            prev_capex + float(fp.capex),
            prev_opex + float(fp.opex),
        )

    fp_rows: list[list[str]] = []
    for year in range(1, 11):
        entry = fp_by_year.get(year)
        if entry is None:
            fp_rows.append([f"Y{year}", "—", "—"])
        else:
            capex_total, opex_total = entry
            fp_rows.append([
                f"Y{year}",
                _fmt_money(capex_total),
                _fmt_money(opex_total),
            ])

    _add_simple_table(
        slide,
        Inches(6.8),
        Inches(1.7),
        Inches(6.0),
        Inches(2.5),
        ["Год", "CAPEX ₽", "OPEX ₽"],
        fp_rows,
    )

    # Phase 8.8: OPEX по категориям маркетингового бюджета
    if opex_by_category:
        _add_text_box(
            slide,
            Inches(6.8),
            Inches(4.4),
            Inches(6.0),
            Inches(0.4),
            "OPEX по категориям маркетинга",
            bold=True,
            size=Pt(12),
        )
        cat_rows: list[list[str]] = []
        # Сортируем по сумме убывания
        for cat, total in sorted(
            opex_by_category.items(), key=lambda x: x[1], reverse=True
        ):
            cat_rows.append([_OPEX_CATEGORY_LABELS.get(cat, cat), _fmt_money(total)])
        _add_simple_table(
            slide,
            Inches(6.8),
            Inches(4.85),
            Inches(6.0),
            Inches(2.0),
            ["Категория", "Total ₽"],
            cat_rows,
        )


# Phase 8.8: лейблы маркетинговых категорий OPEX (синхронизировано с
# OPEX_CATEGORIES в backend/app/schemas/financial_plan.py).
_OPEX_CATEGORY_LABELS: dict[str, str] = {
    "digital": "Digital",
    "ecom": "E-com",
    "ooh": "OOH",
    "pr": "PR",
    "smm": "SMM",
    "design": "Design",
    "research": "Research",
    "posm": "ПОСМ",
    "creative": "Creative",
    "special": "Special",
    "merch": "Merch",
    "tv": "TV",
    "listings": "Листинги",
    "other": "Другое",
}


def _build_slide_risks_and_functions(
    prs: PresentationType, project: Project
) -> None:
    """Слайд 11: риски (слева) + готовность функций (справа)."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _add_title(slide, "10. Риски и готовность функций")

    # Левая колонка: риски
    _add_text_box(
        slide,
        Inches(0.5),
        Inches(1.2),
        Inches(6.0),
        Inches(0.4),
        "Ключевые риски",
        bold=True,
        size=Pt(14),
    )
    risks_raw = project.risks or []
    risk_items: list[str] = []
    for r in risks_raw:
        if isinstance(r, str):
            risk_items.append(r)
        elif isinstance(r, dict) and "text" in r:
            risk_items.append(str(r["text"]))
        else:
            risk_items.append(str(r))
    _add_bullets(
        slide,
        Inches(0.5),
        Inches(1.7),
        Inches(6.0),
        Inches(5.0),
        risk_items,
        placeholder="Риски не описаны",
    )

    # Правая колонка: готовность функций
    _add_text_box(
        slide,
        Inches(6.8),
        Inches(1.2),
        Inches(6.0),
        Inches(0.4),
        "Готовность функций",
        bold=True,
        size=Pt(14),
    )
    func_raw = project.function_readiness or {}
    func_rows: list[list[str]] = []
    if isinstance(func_raw, dict):
        for dept, entry in func_raw.items():
            if not isinstance(entry, dict):
                continue
            status = entry.get("status", "")
            notes = entry.get("notes", "") or ""
            status_label = FUNCTION_STATUS_LABELS.get(str(status), str(status) or "—")
            func_rows.append([str(dept), status_label, notes])

    _add_simple_table(
        slide,
        Inches(6.8),
        Inches(1.7),
        Inches(6.0),
        Inches(5.0),
        ["Функция", "Статус", "Notes"],
        func_rows,
    )


def _build_slide_roadmap_and_approvers(
    prs: PresentationType, project: Project
) -> None:
    """Слайд 12: roadmap (сверху) + согласующие (снизу)."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _add_title(slide, "11. Дорожная карта и согласующие")

    # Roadmap — верхняя половина
    _add_text_box(
        slide,
        Inches(0.5),
        Inches(1.2),
        Inches(12.3),
        Inches(0.4),
        "Дорожная карта",
        bold=True,
        size=Pt(14),
    )
    roadmap_raw = project.roadmap_tasks or []
    roadmap_rows: list[list[str]] = []
    if isinstance(roadmap_raw, list):
        for task in roadmap_raw:
            if not isinstance(task, dict):
                continue
            roadmap_rows.append([
                str(task.get("name", "")),
                str(task.get("start_date", "") or "—"),
                str(task.get("end_date", "") or "—"),
                str(task.get("status", "") or "—"),
                str(task.get("owner", "") or "—"),
            ])
    _add_simple_table(
        slide,
        Inches(0.5),
        Inches(1.7),
        Inches(12.3),
        Inches(2.6),
        ["Задача", "Начало", "Конец", "Статус", "Ответственный"],
        roadmap_rows,
    )

    # Approvers — нижняя половина
    _add_text_box(
        slide,
        Inches(0.5),
        Inches(4.5),
        Inches(12.3),
        Inches(0.4),
        "Согласующие",
        bold=True,
        size=Pt(14),
    )
    approvers_raw = project.approvers or []
    approvers_rows: list[list[str]] = []
    if isinstance(approvers_raw, list):
        for a in approvers_raw:
            if not isinstance(a, dict):
                continue
            # backward compat: в 4.5.1 ELEKTRA-сиде ключ был "approver",
            # в 4.5.3 UI — "name". Поддерживаем оба.
            name = a.get("name") or a.get("approver") or ""
            approvers_rows.append([
                str(a.get("metric", "") or "—"),
                str(name) or "—",
                str(a.get("source", "") or "—"),
            ])
    _add_simple_table(
        slide,
        Inches(0.5),
        Inches(5.0),
        Inches(12.3),
        Inches(2.0),
        ["Метрика", "Согласующий", "Источник"],
        approvers_rows,
    )


def _build_slide_pricing(
    prs: PresentationType,
    pricing: Any | None,
) -> None:
    """Phase 8.1: слайд ценовой сводки SKU × канал."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _add_title(slide, "Цены: полка / ex-factory / COGS")

    if pricing is None or not pricing.skus:
        _add_text_box(
            slide,
            Inches(0.5),
            Inches(2.5),
            Inches(12.0),
            Inches(0.5),
            "Нет SKU с привязанными каналами.",
            size=BODY_FONT_SIZE,
        )
        return

    # Собираем уникальные каналы по всем SKU
    channel_codes: list[str] = []
    seen: set[str] = set()
    for s in pricing.skus:
        for c in s.channels:
            if c.channel_code not in seen:
                seen.add(c.channel_code)
                channel_codes.append(c.channel_code)

    # Header: ["Канал"] + [sku_brand sku_name за каждый SKU]
    headers = ["Канал"] + [f"{s.sku_brand} {s.sku_name}" for s in pricing.skus]

    # Rows: для каждого канала строка [code, shelf_reg, shelf_reg, ...]
    # Покажем 3 секции: shelf, ex-factory, channel margin
    rows: list[list[str]] = []

    # Shelf prices
    rows.append(["— Цена полки, ₽ —"] + [""] * len(pricing.skus))
    for code in channel_codes:
        row = [code]
        for s in pricing.skus:
            cell = next((c for c in s.channels if c.channel_code == code), None)
            row.append(_fmt_money(float(cell.shelf_price_reg), 2) if cell else "—")
        rows.append(row)

    # Ex-Factory
    rows.append(["— Ex-Factory, ₽ —"] + [""] * len(pricing.skus))
    for code in channel_codes:
        row = [code]
        for s in pricing.skus:
            cell = next((c for c in s.channels if c.channel_code == code), None)
            row.append(_fmt_money(float(cell.ex_factory), 2) if cell else "—")
        rows.append(row)

    # COGS row
    cogs_row = ["COGS / шт"]
    for s in pricing.skus:
        cogs_row.append(_fmt_money(float(s.cogs_per_unit), 2))
    rows.append(cogs_row)

    _add_simple_table(
        slide,
        Inches(0.3),
        Inches(1.2),
        Inches(12.7),
        Inches(5.5),
        headers,
        rows,
    )


def _build_slide_value_chain(
    prs: PresentationType,
    value_chain: Any | None,
) -> None:
    """Phase 8.2: слайд per-unit waterfall (Стакан)."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _add_title(slide, "Стакан: per-unit экономика SKU × канал")

    if value_chain is None or not value_chain.skus:
        _add_text_box(
            slide,
            Inches(0.5),
            Inches(2.5),
            Inches(12.0),
            Inches(0.5),
            "Нет SKU с привязанными каналами.",
            size=BODY_FONT_SIZE,
        )
        return

    # Возьмём первый SKU и его каналы для компактного waterfall
    sku = value_chain.skus[0]
    if not sku.channels:
        _add_text_box(
            slide,
            Inches(0.5),
            Inches(2.5),
            Inches(12.0),
            Inches(0.5),
            f"Для SKU «{sku.sku_brand} {sku.sku_name}» нет каналов.",
            size=BODY_FONT_SIZE,
        )
        return

    _add_text_box(
        slide,
        Inches(0.5),
        Inches(1.0),
        Inches(12.0),
        Inches(0.4),
        f"SKU: {sku.sku_brand} {sku.sku_name}"
        + (f" ({sku.sku_volume_l}л)" if sku.sku_volume_l else ""),
        bold=True,
        size=Pt(13),
    )

    headers = ["Показатель"] + [c.channel_code for c in sku.channels]
    waterfall = [
        ("Цена полки", "shelf_price_reg"),
        ("Ex-Factory", "ex_factory"),
        ("COGS итого", "cogs_total"),
        ("Валовая прибыль", "gross_profit"),
        ("Логистика", "logistics"),
        ("Contribution", "contribution"),
        ("CA&M", "ca_m"),
        ("Маркетинг", "marketing"),
        ("EBITDA", "ebitda"),
    ]
    rows: list[list[str]] = []
    for label, attr in waterfall:
        row = [label]
        for cell in sku.channels:
            val = getattr(cell, attr)
            row.append(_fmt_money(float(val), 2))
        rows.append(row)

    # Margins as percentages
    for label, attr in [
        ("GP %", "gp_margin"),
        ("CM %", "cm_margin"),
        ("EBITDA %", "ebitda_margin"),
    ]:
        row = [label]
        for cell in sku.channels:
            row.append(_fmt_pct(float(getattr(cell, attr))))
        rows.append(row)

    _add_simple_table(
        slide,
        Inches(0.3),
        Inches(1.5),
        Inches(12.7),
        Inches(5.2),
        headers,
        rows,
    )


def _build_slide_sensitivity(
    prs: PresentationType,
    sensitivity_data: dict | None,
) -> None:
    """Слайд: 2D sensitivity matrix (Phase 8.4).

    Rows = параметры (ND / Offtake / Shelf Price / COGS).
    Columns = дельты (-20% / -10% / 0% / +10% / +20%).
    Cells = NPV Y1-Y10.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _add_title(slide, "Анализ чувствительности — NPV Y1-Y10")

    if sensitivity_data is None:
        _add_text_box(
            slide,
            Inches(0.5),
            Inches(2.5),
            Inches(12.0),
            Inches(1.0),
            "Нет данных для анализа чувствительности (SKU/каналы не настроены).",
            size=Pt(14),
        )
        return

    base_npv = sensitivity_data.get("base_npv_y1y10")
    deltas = sensitivity_data.get("deltas", [])
    params = sensitivity_data.get("params", [])
    cells = sensitivity_data.get("cells", [])

    # Subtitle with base NPV
    _add_text_box(
        slide,
        Inches(0.5),
        Inches(1.2),
        Inches(12.0),
        Inches(0.4),
        f"Базовый NPV Y1-Y10: {_fmt_money(base_npv)} ₽",
        size=Pt(12),
    )

    # Build lookup: (param, delta) → npv
    npv_lookup: dict[tuple[str, float], float | None] = {}
    for c in cells:
        npv_lookup[(c["parameter"], c["delta"])] = c.get("npv_y1y10")

    param_labels = {
        "nd": "Дистрибуция (ND)",
        "offtake": "Офтейк",
        "shelf_price": "Цена полки",
        "cogs": "Себестоимость (COGS)",
    }

    # Table: header row + 1 row per param
    headers = ["Параметр"] + [f"{d:+.0%}" for d in deltas]
    rows: list[list[str]] = []
    for p in params:
        row = [param_labels.get(p, p)]
        for d in deltas:
            npv = npv_lookup.get((p, d))
            row.append(_fmt_money(npv))
        rows.append(row)

    _add_simple_table(
        slide,
        Inches(0.5),
        Inches(1.8),
        Inches(12.3),
        Inches(4.5),
        headers,
        rows,
    )


def _build_slide_market_and_supply(
    prs: PresentationType, project: Project
) -> None:
    """Phase 8.9 + 8.10: Nielsen бенчмарки (слева) + КП поставщиков (справа)."""
    nielsen = project.nielsen_benchmarks or []
    quotes = project.supplier_quotes or []

    if not nielsen and not quotes:
        return  # Skip slide entirely if no data

    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _add_title(slide, "Рынок и поставки")

    # Левая колонка: Nielsen бенчмарки
    _add_text_box(
        slide,
        Inches(0.5),
        Inches(1.2),
        Inches(6.0),
        Inches(0.4),
        "Nielsen бенчмарки рынка",
        bold=True,
        size=Pt(13),
    )

    if nielsen:
        rows: list[list[str]] = []
        for n in nielsen:
            if not isinstance(n, dict):
                continue
            rows.append([
                str(n.get("channel", "—")),
                _fmt_money(n.get("universe_outlets")),
                (
                    f"{float(n['offtake']):.2f}"
                    if n.get("offtake") is not None
                    else "—"
                ),
                _fmt_pct(n.get("nd_pct")),
                _fmt_money(n.get("avg_price"), 2),
            ])
        _add_simple_table(
            slide,
            Inches(0.5),
            Inches(1.7),
            Inches(6.0),
            Inches(5.0),
            ["Канал", "Universe", "Offtake", "ND %", "Цена ₽"],
            rows,
        )
    else:
        _add_text_box(
            slide,
            Inches(0.5),
            Inches(1.7),
            Inches(6.0),
            Inches(0.5),
            "Нет данных.",
            size=BODY_FONT_SIZE,
        )

    # Правая колонка: КП поставщиков
    _add_text_box(
        slide,
        Inches(6.8),
        Inches(1.2),
        Inches(6.0),
        Inches(0.4),
        "КП на производство",
        bold=True,
        size=Pt(13),
    )

    if quotes:
        rows = []
        for q in quotes:
            if not isinstance(q, dict):
                continue
            rows.append([
                str(q.get("supplier", "—")),
                str(q.get("item", "—")),
                _fmt_money(q.get("price_per_unit"), 2),
                str(q.get("unit") or "—"),
                str(q.get("lead_time_days") or "—"),
            ])
        _add_simple_table(
            slide,
            Inches(6.8),
            Inches(1.7),
            Inches(6.0),
            Inches(5.0),
            ["Поставщик", "Позиция", "Цена ₽", "Ед.", "Срок"],
            rows,
        )
    else:
        _add_text_box(
            slide,
            Inches(6.8),
            Inches(1.7),
            Inches(6.0),
            Inches(0.5),
            "Нет КП.",
            size=BODY_FONT_SIZE,
        )


def _build_slide_executive_summary(
    prs: PresentationType, project: Project
) -> None:
    """Слайд 13: Executive summary.

    Приоритет: AI-generated (Phase 7.4) > manual (Phase 4.5).
    Если оба пусты — fallback placeholder.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BLANK])
    _add_title(slide, "12. Executive Summary")

    # AI-generated имеет приоритет (прошёл review аналитика перед сохранением)
    text = project.ai_executive_summary or project.executive_summary or (
        "Executive summary ещё не сгенерирован."
    )

    _add_text_box(
        slide,
        Inches(0.5),
        Inches(1.5),
        Inches(12.3),
        Inches(5.5),
        text,
        size=Pt(14),
    )


# ============================================================
# Data loading for PPT (media images)
# ============================================================


async def _load_package_images(
    session: AsyncSession,
    skus_with_bom: list[tuple[ProjectSKU, list[BOMItem]]],
) -> dict[int, Path]:
    """Загружает file paths для package_image_id каждого SKU.

    Возвращает {media_id: absolute Path}. Пропускает SKU без image
    или с битыми ссылками — файл просто не окажется на слайде.
    """
    out: dict[int, Path] = {}
    media_ids = [
        psk.package_image_id
        for psk, _ in skus_with_bom
        if psk.package_image_id is not None
    ]
    if not media_ids:
        return out

    assets = (
        await session.scalars(
            select(MediaAsset).where(MediaAsset.id.in_(media_ids))
        )
    ).all()

    storage_root = Path(settings.media_storage_root)
    for asset in assets:
        path = storage_root / asset.storage_path
        if path.is_file():
            out[asset.id] = path
    return out


# ============================================================
# Public entry point
# ============================================================


async def generate_project_pptx(
    session: AsyncSession,
    project_id: int,
) -> bytes:
    """Генерирует PPTX для проекта, возвращает bytes.

    Raises:
        ProjectNotFoundForExport: если project_id не существует.
    """
    project = await _load_project_full(session, project_id)
    if project is None:
        raise ProjectNotFoundForExport(f"Project {project_id} not found")

    # Data loading (аналогично excel_exporter)
    inflation_profile: RefInflation | None = None
    if project.inflation_profile_id is not None:
        inflation_profile = await session.get(
            RefInflation, project.inflation_profile_id
        )

    skus_with_bom = await _load_skus_with_bom(session, project_id)
    psk_channels = await _load_psk_channels(session, project_id)
    package_images = await _load_package_images(session, skus_with_bom)

    fp_rows = (
        await session.scalars(
            select(ProjectFinancialPlan).where(
                ProjectFinancialPlan.project_id == project_id
            )
        )
    ).all()

    # Phase 8.8: OPEX по категориям (aggregated по проекту)
    from app.models import OpexItem
    opex_by_category: dict[str, float] = {}
    if fp_rows:
        fp_ids = [fp.id for fp in fp_rows]
        opex_items = (
            await session.scalars(
                select(OpexItem).where(OpexItem.financial_plan_id.in_(fp_ids))
            )
        ).all()
        for oi in opex_items:
            cat = oi.category or "other"
            opex_by_category[cat] = opex_by_category.get(cat, 0.0) + float(oi.amount)

    sorted_periods, period_by_id = await _load_period_catalog(session)

    scenarios = (
        await session.scalars(
            select(Scenario).where(Scenario.project_id == project_id)
        )
    ).all()
    results_by_scenario = await _load_scenario_results(session, project_id)

    # Base pipeline (для PnL годового)
    base_scenario = next(
        (s for s in scenarios if s.type == ScenarioType.BASE), None
    )
    base_aggregate: Any | None = None
    if base_scenario is not None and skus_with_bom and psk_channels:
        try:
            line_inputs = await build_line_inputs(
                session, project_id, base_scenario.id
            )
            capex, opex = await _load_project_financial_plan(
                session, project_id, sorted_periods
            )
            base_aggregate = run_project_pipeline(
                line_inputs, project_capex=capex, project_opex=opex
            )
        except Exception:  # noqa: BLE001
            base_aggregate = None

    # Build presentation (16:9)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _build_slide_title(prs, project)
    _build_slide_general_info(prs, project)
    _build_slide_concept(prs, project)
    _build_slide_technology(prs, project)
    _build_slide_validation(prs, project)
    _build_slide_sku_mix(prs, skus_with_bom, package_images)
    _build_slide_financial_model(prs, project, inflation_profile)
    _build_slide_kpi(prs, list(scenarios), results_by_scenario)
    _build_slide_pnl(prs, base_aggregate)

    # Pricing + Value Chain (Phase 8.1 / 8.2)
    pricing_data: Any | None = None
    value_chain_data: Any | None = None
    try:
        from app.services.pricing_service import (
            build_pricing_summary,
            build_value_chain,
        )
        pricing_data = await build_pricing_summary(session, project)
        value_chain_data = await build_value_chain(session, project)
    except Exception:  # noqa: BLE001
        pass
    _build_slide_pricing(prs, pricing_data)
    _build_slide_value_chain(prs, value_chain_data)

    # Sensitivity matrix (Phase 8.4)
    sensitivity_data: dict | None = None
    if base_scenario is not None:
        try:
            from app.services.sensitivity_service import compute_sensitivity
            sensitivity_data = await compute_sensitivity(
                session, project_id, base_scenario.id
            )
        except Exception:  # noqa: BLE001
            sensitivity_data = None
    _build_slide_sensitivity(prs, sensitivity_data)

    _build_slide_cogs_and_fin_plan(
        prs, skus_with_bom, list(fp_rows), period_by_id,
        opex_by_category=opex_by_category,
    )
    _build_slide_risks_and_functions(prs, project)
    _build_slide_roadmap_and_approvers(prs, project)
    _build_slide_market_and_supply(prs, project)  # Phase 8.9 + 8.10
    _build_slide_executive_summary(prs, project)

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
